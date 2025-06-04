from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# Create a new PowerPoint presentation
presentation = Presentation()

# Add a blank slide (layout index 6 is typically a blank layout)
slide = presentation.slides.add_slide(presentation.slide_layouts[6])

# Add a shape with the specified position and size
x = Inches(0.3125)
y = Inches(4.451388888888889)
width = Inches(0.625)
height = Inches(0.4652777777777778)

# Create a rectangle shape
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,  # MSO_SHAPE.RECTANGLE corresponds to a rectangle shape
    left=x,
    top=y,
    width=width,
    height=height
)

# Set the background color of the shape to gray
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(169, 169, 169)  # Light gray color

# Add text to the shape
text_frame = shape.text_frame
text_frame.text = "2016"

# Format the text
paragraph = text_frame.paragraphs[0]
run = paragraph.add_run()
run.text = "2016"
run.font.size = Pt(12)  # Set font size
run.font.name = "Arial"  # Set font name
run.font.underline = True  # Underline the text
run.font.color.rgb = RGBColor(0, 0, 0)  # Black text color

# Center the text horizontally and vertically
text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Vertically center the text
paragraph.alignment = PP_ALIGN.CENTER  # Horizontally center the text

# Ensure the output directory exists
output_dir = "output/generated_ppts"
os.makedirs(output_dir, exist_ok=True)

# Save the presentation
output_path = os.path.join(output_dir, "generated_slide.pptx")
presentation.save(output_path)

print(f"Presentation saved to {output_path}")