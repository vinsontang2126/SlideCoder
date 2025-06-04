from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Create a new PowerPoint presentation
presentation = Presentation()

# Add a blank slide (layout index 6 is typically a blank layout)
slide = presentation.slides.add_slide(presentation.slide_layouts[6])

# Add the textual content
text_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(4))
text_frame = text_box.text_frame
text_frame.word_wrap = True

# Add the title
title_paragraph = text_frame.add_paragraph()
title_paragraph.text = "Analysis of the Region:"
title_paragraph.font.bold = True
title_paragraph.font.size = Pt(24)
title_paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black color
title_paragraph.alignment = PP_ALIGN.LEFT

# Add the textual content
content_paragraph = text_frame.add_paragraph()
content_paragraph.text = (
    "Textual Content\n"
    "The textual content provides an introduction to the concept of text frames within shapes in PowerPoint (PPT). "
    "It explains the functionality and properties of text frames, emphasizing their role in managing text within shapes. "
    "Key points include:\n"
    "1. Text Frame Properties:\n"
    "   - Vertical alignment.\n"
    "   - Margins.\n"
    "   - Text wrapping.\n"
    "   - Auto-fit behavior.\n"
    "   - Rotation angle.\n"
    "   - Optional 3D visual effects.\n"
    "   - Multi-column formatting.\n\n"
    "2. Paragraph Structure:\n"
    "   - The text frame always contains a sequence of paragraphs.\n"
    "   - At least one paragraph is present, even if it is empty."
)
content_paragraph.font.size = Pt(14)
content_paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black color
content_paragraph.alignment = PP_ALIGN.LEFT

# Add the graphical element (gray background with red lines)
graphic_shape = slide.shapes.add_shape(
    # Rectangle shape for the gray background
    1,  # MSO_SHAPE.RECTANGLE
    Inches(5.3125), Inches(4.6875), Inches(4.0625), Inches(1.1666666666666667)
)
graphic_shape.fill.solid()
graphic_shape.fill.fore_color.rgb = RGBColor(192, 192, 192)  # Gray background

# Add red lines to represent paragraphs
line_positions = [0.1, 0.3, 0.5, 0.7, 0.9]  # Relative positions within the gray background
line_widths = [3.5, 3.0, 2.5, 2.0, 1.5]  # Uneven line widths
for i, (pos, width) in enumerate(zip(line_positions, line_widths)):
    line = slide.shapes.add_shape(
        1,  # MSO_SHAPE.LINE
        Inches(5.3125 + 0.1), Inches(4.6875 + pos), Inches(width), Inches(0.05)
    )
    line.line.color.rgb = RGBColor(255, 0, 0)  # Red color
    line.line.width = Pt(2)  # Line thickness

# Save the presentation
output_dir = "output/generated_ppts"
os.makedirs(output_dir, exist_ok=True)  # Create the directory if it doesn't exist
output_path = os.path.join(output_dir, "generated_slide.pptx")
presentation.save(output_path)

print(f"Presentation saved at: {output_path}")