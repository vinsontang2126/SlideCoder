from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Create a new PowerPoint presentation
presentation = Presentation()

# Add a blank slide (layout index 6 is typically a blank layout)
slide = presentation.slides.add_slide(presentation.slide_layouts[6])

# Add the textual content box
text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
text_frame = text_box.text_frame

# Add the title
title_paragraph = text_frame.add_paragraph()
title_paragraph.text = "Detailed Analysis of the Region"
title_paragraph.font.bold = True
title_paragraph.font.size = Pt(24)
title_paragraph.font.color.rgb = RGBColor(0, 0, 0)
title_paragraph.alignment = PP_ALIGN.LEFT

# Add the textual content
content_paragraph = text_frame.add_paragraph()
content_paragraph.text = (
    "Textual Content\n\n"
    "The textual content provides an introduction to the concept of text frames within shapes in PowerPoint (PPT). "
    "It explains the functionality and properties of text frames, emphasizing their role in managing text within shapes. "
    "Key points include:\n"
    "1. Text Frame Properties:\n"
    "   - Vertical alignment.\n"
    "   - Margins.\n"
    "   - Text wrapping.\n"
    "   - Auto-fit behavior.\n"
    "   - Rotation angle.\n"
    "   - Optional 3D visual effects.\n"
    "   - Multi-column formatting.\n\n"
    "2. Paragraph Structure:\n"
    "   - The text frame always contains a sequence of paragraphs.\n"
    "   - At least one paragraph is present, even if it is empty.\n\n"
    "The content is technical and descriptive, aimed at users who want to understand the mechanics of text handling in PowerPoint shapes."
)
content_paragraph.font.size = Pt(14)
content_paragraph.font.color.rgb = RGBColor(0, 0, 0)
content_paragraph.alignment = PP_ALIGN.LEFT

# Add the graphical element (rectangle shape with gray background)
graphical_element = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,  # Rectangle shape
    Inches(8.125),
    Inches(7.027777777777778),
    Inches(1.5625),
    Inches(0.47222222222222215)
)
graphical_element.fill.solid()
graphical_element.fill.fore_color.rgb = RGBColor(169, 169, 169)  # Gray background

# Add text to the graphical element
graphical_text_frame = graphical_element.text_frame
graphical_paragraph = graphical_text_frame.add_paragraph()
graphical_paragraph.text = "March 10, 2018"
graphical_paragraph.font.size = Pt(12)
graphical_paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black text
graphical_paragraph.font.underline = True
graphical_paragraph.alignment = PP_ALIGN.CENTER

# Ensure the output directory exists
output_dir = "output/generated_ppts"
os.makedirs(output_dir, exist_ok=True)

# Save the presentation
output_path = os.path.join(output_dir, "generated_slide.pptx")
presentation.save(output_path)

print(f"Presentation saved to {output_path}")