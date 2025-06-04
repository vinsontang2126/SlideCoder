from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_VERTICAL_ANCHOR
import os

# Create a new PowerPoint presentation
presentation = Presentation()

# Add a blank slide (layout index 6 is typically a blank layout)
slide = presentation.slides.add_slide(presentation.slide_layouts[6])

# Set the slide background to solid gray
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(169, 169, 169)  # Gray color (RGB: 169, 169, 169)

# Add the graphical text ("...if not, Sci-Hub would not exist")
text_box = slide.shapes.add_textbox(Inches(0.3125), Inches(0.22916666666666666), Inches(8.125), Inches(0.7083333333333333))
text_frame = text_box.text_frame
text_frame.word_wrap = True  # Enable text wrapping

# Add the text content
p = text_frame.add_paragraph()
p.text = "...if not, Sci-Hub would not exist"
p.font.bold = True
p.font.size = Pt(32)  # Font size
p.font.name = "Arial"  # Sans-serif font
p.font.color.rgb = RGBColor(0, 0, 0)  # Black color

# Correctly set vertical alignment to middle
text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE  # Vertically center the text

# Ensure the output directory exists
output_dir = "output/generated_ppts"
os.makedirs(output_dir, exist_ok=True)

# Save the presentation
output_path = os.path.join(output_dir, "generated_slide.pptx")
presentation.save(output_path)

print(f"Presentation saved at: {output_path}")