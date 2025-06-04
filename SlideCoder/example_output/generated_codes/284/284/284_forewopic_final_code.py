from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Create a new PowerPoint presentation
presentation = Presentation()

# Add a blank slide (layout index 6 is typically a blank layout)
slide = presentation.slides.add_slide(presentation.slide_layouts[6])

# Set the slide background to solid gray
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(169, 169, 169)  # Gray color (RGB: 169, 169, 169)

# --- Top Section: Add the main text ---
top_text_box = slide.shapes.add_textbox(Inches(0.3125), Inches(0.22916666666666666), Inches(8.125), Inches(0.7083333333333333))
top_text_frame = top_text_box.text_frame
top_text_frame.word_wrap = True  # Enable text wrapping

# Add the text content
top_paragraph = top_text_frame.add_paragraph()
top_paragraph.text = "...if not, Sci-Hub would not exist"
top_paragraph.font.bold = True
top_paragraph.font.size = Pt(32)  # Font size
top_paragraph.font.name = "Arial"  # Sans-serif font
top_paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black color
top_text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Vertically center the text

# --- Middle Section: Add the year and red horizontal lines ---
# Add the year "2016"
year_text_box = slide.shapes.add_textbox(Inches(0.3125), Inches(1.5), Inches(1), Inches(0.5))
year_text_frame = year_text_box.text_frame
year_text_frame.word_wrap = True

year_paragraph = year_text_frame.add_paragraph()
year_paragraph.text = "2016"
year_paragraph.font.size = Pt(12)  # Font size
year_paragraph.font.name = "Arial"  # Font name
year_paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black color
year_paragraph.alignment = PP_ALIGN.LEFT  # Left-aligned

# Add red horizontal lines
line_positions = [1.8, 2.0, 2.2, 2.4, 2.6]  # Vertical positions for the lines
line_widths = [3.5, 3.0, 2.5, 2.0, 1.5]  # Line widths in inches
for i, (pos, width) in enumerate(zip(line_positions, line_widths)):
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,  # Corrected: Rectangle shape for the line
        Inches(5), Inches(pos), Inches(width), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red color

# --- Bottom Section: Add the date ---
date_text_box = slide.shapes.add_textbox(Inches(7.5), Inches(5.5), Inches(2), Inches(0.5))
date_text_frame = date_text_box.text_frame
date_text_frame.word_wrap = True

date_paragraph = date_text_frame.add_paragraph()
date_paragraph.text = "March 10, 2018"
date_paragraph.font.size = Pt(12)  # Font size
date_paragraph.font.name = "Arial"  # Font name
date_paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black color
date_paragraph.alignment = PP_ALIGN.RIGHT  # Right-aligned

# --- Save the presentation ---
# Ensure the output directory exists
output_dir = "output/generated_ppts"
os.makedirs(output_dir, exist_ok=True)

# Save the presentation
output_path = os.path.join(output_dir, "generated_slide.pptx")
presentation.save(output_path)

print(f"Presentation saved to {output_path}")